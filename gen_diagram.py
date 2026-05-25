#!/usr/bin/env python3
"""
DDR3 乒乓缓冲 — 正交走线框图 → .pptx (Visio 可直接导入)
"""
import math, os

OUT_DIR = r'E:\work_space\Claude_Work\DDR_pingpang'

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE

# ═══════════════ 配色 ═══════════════
WR= '1565C0'; RD= '2E7D32'; CTL= 'E65100'; XCK= 'C62828'
GRAY= '546E7A'; DARK= '212121'; LITE= '616161'
GREEN=RD; ORANGE=CTL; RED=XCK; BLUE=WR

def rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# ═══════════════ 布局 (英寸) ═══════════════
# 列 X (模块左上角)
CX = {
    'port':   0.30,   # 输入端口
    'ctrl':   1.50,   # 写控制 / 读控制
    'dist':   2.90,   # 数据分发 / 输出MUX
    'fifo_w': 4.40,   # 写 FIFO
    'fifo_r': 4.40,   # 读 FIFO
    'fsm':    6.00,   # rd_flag 状态机
    'mig':    7.40,   # MIG
    'ddr':    9.80,   # DDR3 颗粒
    'out_fifo': 2.90, # 输出 FIFO
    'dout':   4.40,   # 输出端口
    'tb':     6.20,   # TB 校验
    'text':   10.80,  # 文字说明
}

# 模块宽高 (英寸)
W = {
    'port': 1.00, 'ctrl': 1.20, 'dist': 1.30, 'fifo': 1.40,
    'fsm': 1.20, 'mig': 2.20, 'ddr': 1.50, 'out_fifo': 1.30,
    'dout': 1.20, 'tb': 2.00,
}
H = {
    'port': 0.55, 'ctrl': 0.95, 'dist': 0.85, 'fifo': 0.55,
    'fsm': 0.65, 'mig': 1.50, 'ddr': 1.20, 'out_fifo': 0.90,
    'dout': 0.90, 'tb': 0.90,
}

# ═══════════════ 辅助函数 ═══════════════
def box(slide, x, y, w, h, label, fill, stroke, sz=10):
    left, top = Inches(x), Inches(y)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
    s.line.color.rgb = rgb(stroke); s.line.width = Pt(1.2)
    tf = s.text_frame; tf.word_wrap = True
    for i, ln in enumerate(label.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(sz); p.font.color.rgb = rgb(DARK)
        p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    return s

def txt(slide, x, y, text, clr=DARK, sz=8, bold=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.5), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = rgb(clr); p.font.bold = bold
    return tb

def arrow(slide, x1, y1, x2, y2, clr, sw=2.0, dash=False):
    """ELBOW 正交走线 + 三角箭头"""
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.ELBOW,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = rgb(clr); c.line.width = Pt(sw)
    if dash: c.line.dash_style = 2

    # 箭头 XML
    ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    spPr = c._element.find(f'{ns_p}spPr')
    if spPr is None:
        spPr = c._element.makeelement(f'{ns_p}spPr', {})
        c._element.insert(0, spPr)
    ln = spPr.find(f'{ns_a}ln')
    if ln is None:
        ln = spPr.makeelement(f'{ns_a}ln', {})
        spPr.append(ln)
    tail = ln.find(f'{ns_a}tailEnd')
    if tail is None:
        tail = ln.makeelement(f'{ns_a}tailEnd', {})
        ln.append(tail)
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return c

def channel_bg(slide, x, y, w, h, clr):
    """通道背景色块"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(clr)
    s.line.fill.background()
    return s

# ═══════════════ 绘制 ═══════════════
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    txt(slide, 4.5, 0.05, 'DDR3 Ping-Pong Buffer  —  FPGA Block Diagram', '0D47A1', 16, True)

    # ═══════════ 通道 0 ═══════════
    Y0 = 0.55
    channel_bg(slide, 0.15, Y0, 12.9, 2.45, 'FFF3E0')
    txt(slide, 0.25, Y0+0.05, 'Channel 0 — fifo_ddr_0  (mig_7series_0,  ADDR_INITIAL=0x0000_0000)', 'E65100', 10, True)

    # 第1行: 输入/端口 (y=Y0+0.3)
    y1 = Y0 + 0.30
    box(slide, CX['port'], y1, W['port'], H['port'], 'din[31:0]', 'E8F5E9', GREEN)
    txt(slide, CX['port']+0.15, y1+H['port']+0.02, 'din_empty/alempty', LITE, 7)
    box(slide, CX['ctrl'], y1, W['ctrl'], H['ctrl'], 'Write Control\nwr_flag_0 / wr_cnt', 'FFFDE7', 'F9A825', 9)
    box(slide, CX['dist'], y1, W['dist'], H['dist'], 'Data Distribute\nddr3_0_din = din', 'E8EAF6', '5C6BC0', 9)
    box(slide, CX['fifo_w'], y1, W['fifo'], H['fifo'], 'Write FIFO\nasync_1024x32', 'F3E5F5', '7B1FA2', 9)
    box(slide, CX['fsm'], y1, W['fsm'], H['fsm'], 'rd_flag FSM\n+ Addr Gen', 'FFFDE7', 'F9A825', 9)
    # MIG 从 y1-0.05 开始（稍高）
    box(slide, CX['mig'], y1-0.10, W['mig'], H['mig'], 'mig_7series_0\n\nMIG IP  (64-bit app I/F)\napp_cmd: 000=W / 001=R\napp_wdf_data = {32\'d0, dout}\napp_en = en_wr | en_rd', 'FFEBEE', RED, 8)
    box(slide, CX['ddr'], y1, W['ddr'], H['ddr'], 'DDR3 SDRAM\nParticle #0\n16-bit data bus', 'ECEFF1', GRAY, 9)

    # 第2行: 读 FIFO (y = y1+H['mig']+0.15)
    y1b = y1 + H['mig'] + 0.15
    box(slide, CX['fifo_r'], y1b, W['fifo'], H['fifo'], 'Read FIFO\nasync_1024x32', 'F3E5F5', '7B1FA2', 9)

    # ── 通道0 箭头 ──
    # 写数据流 (L→R, 蓝色粗线)
    arrow(slide, CX['port']+W['port'], y1+H['port']/2,  CX['dist'],       y1+H['dist']/2,   WR, 2.5)  # din→dist
    arrow(slide, CX['dist']+W['dist'],     y1+H['dist']/2,  CX['fifo_w'],    y1+H['fifo']/2,   WR, 2.5)  # dist→wr_fifo
    arrow(slide, CX['fifo_w']+W['fifo'],   y1+H['fifo']/2,  CX['mig'],       y1+0.35,           WR, 2.5)  # wr_fifo→MIG(wr)
    arrow(slide, CX['mig']+W['mig'],       y1+H['mig']/2-0.1, CX['ddr'],      y1+H['ddr']/2,   XCK, 2.0, True)  # MIG→DDR3 (PHY)

    # 控制信号 (橙色虚线)
    arrow(slide, CX['ctrl'],             y1+H['ctrl']*0.65, CX['dist'],      y1+H['dist']*0.65, CTL, 1.5, True)  # wr_ctrl→dist
    arrow(slide, CX['fsm']+W['fsm'],     y1+H['fsm']/2,     CX['mig']+0.3,  y1+0.65,           CTL, 1.5, True)  # rd_flag→MIG

    # 读数据流 (R→L, 绿色粗线)
    arrow(slide, CX['ddr'],              y1+H['ddr']-0.25,  CX['mig']+W['mig'], y1+H['mig']-0.30, RD, 2.0)  # DDR3→MIG(read)
    arrow(slide, CX['mig'],              y1b+H['fifo']/2,   CX['fifo_r']+W['fifo'], y1b+H['fifo']/2, RD, 2.5)  # MIG→rd_fifo(bottom)
    # 读FIFO 输出向下
    arrow(slide, CX['fifo_r']+W['fifo']/2, y1b+H['fifo'],  CX['fifo_r']+W['fifo']/2, y1b+H['fifo']+0.25, RD, 2.5)

    # 跨时钟域 (红色虚线)
    arrow(slide, CX['ctrl']+W['ctrl']*0.5, y1+H['ctrl'],   CX['fsm']+W['fsm']*0.3,  y1,               XCK, 1.5, True)  # wr_flag→clk_ddr3
    arrow(slide, CX['mig']+0.2,            y1+H['mig']-0.1, CX['ctrl']+W['ctrl']*0.7, y1+H['ctrl'],  XCK, 1.5, True)  # wr_full→clk

    # 背压 (橙色反向)
    arrow(slide, CX['fifo_w'],             y1+H['fifo']-0.1, CX['dist']+W['dist'],  y1+H['dist']-0.12, CTL, 1.5, True)

    # 信号标注
    txt(slide, CX['port']+W['port']+0.05, y1-0.05, 'din[31:0]', WR, 7)
    txt(slide, CX['dist']-0.30, y1+0.05, 'wr_flag_0', CTL, 6)
    txt(slide, CX['dist']+W['dist']+0.02, y1-0.05, 'ddr3_0_din[31:0]', WR, 7)
    txt(slide, CX['fifo_w']+W['fifo']+0.02, y1-0.10, '→ app_wdf_data[63:0]', WR, 6)
    txt(slide, CX['ddr']-0.20, y1+H['ddr']-0.30, 'app_rd_data[63:0]', RD, 7)
    txt(slide, CX['ctrl']+W['ctrl']+0.02, y1+H['ctrl']-0.05, 'wr_flag_0→\ndata_cross', XCK, 6)
    txt(slide, CX['dist']-0.40, y1+H['dist']-0.05, 'din_prog_full', CTL, 6)

    # ═══════════ 通道 1 ═══════════
    Y1 = 3.15
    channel_bg(slide, 0.15, Y1, 12.9, 2.30, 'E8F5E9')
    txt(slide, 0.25, Y1+0.05, 'Channel 1 — fifo_ddr_1  (mig_7series_1,  ADDR_INITIAL=0xFFFF_FFFF)', '2E7D32', 10, True)

    y2 = Y1 + 0.30
    box(slide, CX['dist'], y2, W['dist'], H['dist'], 'Data Distribute\nddr3_1_din = din\nwr_flag_1=~wr_flag_0', 'E8EAF6', '5C6BC0', 8)
    box(slide, CX['fifo_w'], y2, W['fifo'], H['fifo'], 'Write FIFO\nasync_1024x32', 'F3E5F5', '7B1FA2', 9)
    box(slide, CX['fsm'], y2, W['fsm'], H['fsm'], 'rd_flag FSM\n+ Addr Gen', 'FFFDE7', 'F9A825', 9)
    box(slide, CX['mig'], y2-0.05, W['mig'], 1.30, 'mig_7series_1\nMIG IP (same as ch0)', 'FFEBEE', RED, 8)
    box(slide, CX['ddr'], y2, W['ddr'], H['ddr'], 'DDR3 SDRAM\nParticle #1\n16-bit data bus', 'ECEFF1', GRAY, 9)
    y2b = y2 + H['mig'] + 0.10
    box(slide, CX['fifo_r'], y2b, W['fifo'], H['fifo'], 'Read FIFO\nasync_1024x32', 'F3E5F5', '7B1FA2', 9)

    # 通道1 箭头
    arrow(slide, CX['dist']+W['dist'],   y2+H['dist']/2,  CX['fifo_w'],  y2+H['fifo']/2,  WR, 2.5)
    arrow(slide, CX['fifo_w']+W['fifo'], y2+H['fifo']/2,  CX['mig'],     y2+0.25,          WR, 2.5)
    arrow(slide, CX['mig']+W['mig'],     y2+H['mig']/2-0.05, CX['ddr'],  y2+H['ddr']/2,   XCK, 2.0, True)
    arrow(slide, CX['fsm']+W['fsm'],     y2+H['fsm']/2,    CX['mig']+0.3,y2+0.50,         CTL, 1.5, True)
    arrow(slide, CX['ddr'],              y2+H['ddr']-0.25, CX['mig']+W['mig'], y2+H['mig']-0.25, RD, 2.0)
    arrow(slide, CX['mig'],              y2b+H['fifo']/2,  CX['fifo_r']+W['fifo'], y2b+H['fifo']/2, RD, 2.5)
    arrow(slide, CX['fifo_r']+W['fifo']/2, y2b+H['fifo'], CX['fifo_r']+W['fifo']/2, y2b+H['fifo']+0.20, RD, 2.5)
    arrow(slide, CX['fifo_w'],           y2+H['fifo']-0.1, CX['dist']+W['dist'], y2+H['dist']-0.12, CTL, 1.5, True)
    arrow(slide, CX['ctrl']+W['ctrl']*0.5, y1+H['ctrl']+0.1, CX['fsm']+W['fsm']*0.3, y2, XCK, 1.5, True)
    arrow(slide, CX['mig']+0.2,          y2+1.20,         CX['ctrl']+W['ctrl']*0.7, y1+H['ctrl']+0.15, XCK, 1.5, True)

    txt(slide, CX['dist']+W['dist']+0.02, y2-0.05, 'ddr3_1_din[31:0]', WR, 7)
    txt(slide, CX['fifo_w']+W['fifo']+0.02, y2-0.05, '→ MIG', WR, 6)
    txt(slide, CX['ctrl']+1.0, y1+H['ctrl']+0.15, 'wr_flag_1 = ~wr_flag_0', CTL, 7)

    # ═══════════ 输出路径 ═══════════
    Y2 = 5.65
    channel_bg(slide, 0.15, Y2, 12.9, 1.70, 'E8F0FE')
    txt(slide, 0.25, Y2+0.05, 'Output Path  —  Read-side MUX + Output FIFO + TB Check', '1565C0', 10, True)

    y3 = Y2 + 0.35
    box(slide, CX['dist'], y3, W['dist'], H['out_fifo'], 'Output MUX\nrd_flag_0 selects:\nddr3_0_dout or ddr3_1_dout', 'FCE4EC', 'AD1457', 8)
    box(slide, CX['out_fifo']+1.6, y3, 1.20, H['out_fifo'], 'Output FIFO\nasync_16x32', 'F3E5F5', '7B1FA2', 9)
    box(slide, CX['dout']+1.6, y3, 1.30, H['out_fifo'], 'Read Control\nrd_flag_0 / rd_cnt', 'FFFDE7', 'F9A825', 9)
    box(slide, CX['tb'], y3, W['tb'], H['out_fifo'], 'dout[31:0]\ndout_empty/alempty', 'E8F5E9', GREEN, 9)
    box(slide, CX['text']-0.5, y3, 1.60, H['out_fifo'], 'TB Verify\nddr3_model x2\ndout==ddr_dout+1?', 'E1F5FE', '0277BD', 8)

    # 输出箭头
    # 读FIFO 向下 → MUX
    arrow(slide, CX['fifo_r']+W['fifo']/2, y1b+H['fifo']+0.25, CX['dist']+W['dist']*0.4, y3, RD, 2.5)
    arrow(slide, CX['fifo_r']+W['fifo']/2, y2b+H['fifo']+0.20, CX['dist']+W['dist']*0.7, y3, RD, 2.5)
    # MUX → OutFIFO
    arrow(slide, CX['dist']+W['dist'],     y3+H['out_fifo']/2, CX['out_fifo']+1.6,      y3+H['out_fifo']/2, RD, 2.5)
    # OutFIFO → ReadCtrl → dout
    arrow(slide, CX['out_fifo']+2.8,       y3+H['out_fifo']/2, CX['dout']+1.6,           y3+H['out_fifo']/2, RD, 2.5)
    arrow(slide, CX['dout']+2.9,           y3+H['out_fifo']/2, CX['tb'],                y3+H['out_fifo']/2, RD, 2.5)
    # 控制
    arrow(slide, CX['dout']+1.6,           y3+0.2,             CX['dist']+W['dist']*0.5, y3+0.2, CTL, 1.5, True)
    arrow(slide, CX['out_fifo']+1.6,       y3+H['out_fifo']-0.1, CX['dout']+1.6,        y3+H['out_fifo']-0.1, CTL, 1.5, True)

    txt(slide, CX['dist']+0.8, y3-0.15, 'rd_flag_0', CTL, 7)
    txt(slide, CX['out_fifo']+1.8, y3+H['out_fifo']+0.02, 'fifo_prog_full', CTL, 6)
    txt(slide, CX['dist']+W['dist']+0.02, y3+H['out_fifo']/2-0.05, 'dd3_dout[31:0]', RD, 7)
    txt(slide, CX['dist']+0.1, y3-0.15, 'ddr3_0/1_dout[31:0]', RD, 7)

    # ═══════════ 右侧文字 ═══════════
    tx = CX['text']
    txt(slide, tx, 0.65, 'Ping-Pong Switching:', 'E65100', 10, True)
    for i, s in enumerate(['Phase 1: wr_flag_0=1',
                           '  Write→DDR3_0  Read←DDR3_1',
                           'Phase 2: wr_flag_0=0',
                           '  Write→DDR3_1  Read←DDR3_0',
                           'Toggle: wr_cnt>=addr_size-1',
                           '  && din_rd_en==1']):
        txt(slide, tx, 0.90+i*0.17, s, GRAY if i%2==0 else LITE, 7)

    txt(slide, tx, 3.25, 'Cross-Clock Domain:', XCK, 10, True)
    for i, s in enumerate(['Data path: async FIFO',
                           '  clk <-> clk_ddr3_X',
                           'Reset: reset_cross',
                           '1-bit flag: data_cross',
                           '  wr_flag_X -> clk_ddr3_X',
                           'Output FIFO: same clk']):
        txt(slide, tx, 3.50+i*0.17, s, GRAY if i%2==0 else LITE, 7)

    txt(slide, tx, 4.80, 'Key Signals:', DARK, 10, True)
    for i, s in enumerate(['din/dout: 32-bit user data',
                           'MIG app: 64-bit (upper 0)',
                           'addr_size: 512K words',
                           'wr_flag_0/1: write dir',
                           'rd_flag_0: read MUX select',
                           'prog_full: backpressure']):
        txt(slide, tx, 5.05+i*0.17, s, LITE, 7)

    # 图例
    txt(slide, tx, 6.40, 'Legend:', DARK, 10, True)
    items = [('F3E5F5','7B1FA2','Async FIFO'),('FFEBEE',RED,'MIG IP'),('ECEFF1',GRAY,'DDR3 SDRAM'),
             ('FFFDE7','F9A825','Control'),('FCE4EC','AD1457','MUX/Port'),
             (WR,WR,'Write Data'),(RD,RD,'Read Data'),(CTL,CTL,'Control Sig'),(XCK,XCK,'Cross-CDC')]
    for i,(bg,st,lb) in enumerate(items):
        r = i // 3; c = i % 3
        rx, ry = tx+c*1.10, 6.60+r*0.25
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rx), Inches(ry), Inches(0.18), Inches(0.16))
        b.fill.solid(); b.fill.fore_color.rgb = rgb(bg); b.line.fill.background()
        txt(slide, rx+0.22, ry+0.01, lb, DARK, 7)

    # 保存
    path = os.path.join(OUT_DIR, 'ddr_pingpang_block.pptx')
    prs.save(path)
    print(f'OK: {path}')

if __name__ == '__main__':
    build()
